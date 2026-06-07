#!/usr/bin/env python
"""Build a fully EDITABLE PowerPoint (native text boxes + tables + images)
for the term-project proposal. Marp's --pptx bakes slides into images;
this keeps every title/bullet/table editable in PowerPoint.

Output: docs/slides_proposal_editable.pptx
"""
import os
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
DOCS = os.path.abspath(os.path.join(HERE, ".."))
FIGS = os.path.join(DOCS, "figs")

NAVY = RGBColor(0x1A, 0x36, 0x5D)
SLATE = RGBColor(0x2C, 0x3E, 0x50)
RED = RGBColor(0xC0, 0x39, 0x2B)
DRED = RGBColor(0xA9, 0x32, 0x26)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
REDBG = RGBColor(0xFB, 0xEC, 0xEB)
HEADBG = RGBColor(0x1A, 0x36, 0x5D)
ROWBG = RGBColor(0xF2, 0xF5, 0xF9)
FONT = "Arial"

SW, SH = 13.333, 7.5
prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def textbox(s, text, left, top, width, height, size=18, color=SLATE,
            bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            line_spacing=1.05):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        runs = ln if isinstance(ln, list) else [(ln, {})]
        for j, (t, st) in enumerate(runs):
            r = p.add_run()
            r.text = t
            r.font.name = st.get("font", FONT)
            r.font.size = Pt(st.get("size", size))
            r.font.bold = st.get("bold", bold)
            r.font.italic = st.get("italic", italic)
            r.font.color.rgb = st.get("color", color)
    return tb


def title(s, text, sub=False):
    textbox(s, text, 0.55, 0.3, 12.2, 0.95, size=30, color=NAVY, bold=True,
            anchor=MSO_ANCHOR.MIDDLE)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.22),
                            Inches(12.1), Pt(2.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    ln.line.fill.background()
    return 1.45


def pagenum(s, n):
    textbox(s, str(n), 12.7, 7.0, 0.5, 0.35, size=12, color=GRAY, align=PP_ALIGN.RIGHT)


def place_image(s, name, top, max_w=12.0, max_h=4.6, center_x=None):
    path = os.path.join(FIGS, name)
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w; h = w / ar
    if h > max_h:
        h = max_h; w = h * ar
    left = (SW - w) / 2 if center_x is None else center_x
    s.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))
    return top + h  # bottom


def callout(s, text, left, top, width, height):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                             Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = REDBG
    box.line.color.rgb = RED; box.line.width = Pt(1.5)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    for t, st in text:
        r = p.add_run(); r.text = t
        r.font.name = FONT; r.font.size = Pt(st.get("size", 17))
        r.font.bold = st.get("bold", False); r.font.color.rgb = st.get("color", SLATE)
    return box


# ---- native-shape helpers (used to build the slide-2 diagram editably) ----
GREEN = RGBColor(0x27, 0x86, 0x5A)
BLUE = RGBColor(0x2C, 0x6F, 0xBB)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
THBG = RGBColor(0xEE, 0xF4, 0xFB)
PRBG = RGBColor(0xFD, 0xF0, 0xEF)
LGRAY = RGBColor(0x88, 0x88, 0x88)


def rrect(s, x, y, w, h, fill=WHITE, line=SLATE, line_w=1.5):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def set_text(sp, runs, size=12, align=PP_ALIGN.CENTER, color=SLATE, bold=False):
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("left", "right"):
        setattr(tf, f"margin_{m}", Inches(0.05))
    for m in ("top", "bottom"):
        setattr(tf, f"margin_{m}", Inches(0.02))
    lines = runs if isinstance(runs, list) else [runs]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        rr = ln if isinstance(ln, list) else [(ln, {})]
        for t, st in rr:
            r = p.add_run(); r.text = t
            r.font.name = FONT; r.font.size = Pt(st.get("size", size))
            r.font.bold = st.get("bold", bold); r.font.color.rgb = st.get("color", color)
    return sp


def arrow(s, x1, y1, x2, y2, color=LGRAY, w=1.75, dashed=False):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w)
    ln = cn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    if dashed:
        d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.insert(list(ln).index(tail), d)
    return cn


def arc_curve(s, x0, x1, ybase, peak, color, w=2.5, n=14):
    """Smooth dome drawn as n straight connectors (renders reliably everywhere)."""
    t = np.linspace(0, 1, n + 1)
    xs = x0 + (x1 - x0) * t
    ys = ybase - peak * np.sin(np.pi * t)
    for i in range(n):
        cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(xs[i]), Inches(ys[i]),
                                    Inches(xs[i + 1]), Inches(ys[i + 1]))
        cn.line.color.rgb = color; cn.line.width = Pt(w)


def oval(s, cx, cy, d, color):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2),
                            Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def bar(s, x, y, w, h, color):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def gripper(s, cx, cytop):
    bar(s, cx - 0.28, cytop, 0.56, 0.14, SLATE)            # palm
    bar(s, cx - 0.28, cytop - 0.32, 0.13, 0.34, SLATE)     # left finger
    bar(s, cx + 0.15, cytop - 0.32, 0.13, 0.34, SLATE)     # right finger


def cup_shape(s, cx, ytop, w, h, color):
    sp = s.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(cx - w / 2), Inches(ytop),
                            Inches(w), Inches(h))
    sp.rotation = 180  # wide side up -> cup opening on top
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xF7)
    sp.line.color.rgb = color; sp.line.width = Pt(1.75); sp.shadow.inherit = False
    return sp


def rot_arc(s, cx, cy, r, color, w=2.0):
    """Small rotation arc (alpha) drawn as connectors with an end arrowhead."""
    angs = np.linspace(np.radians(205), np.radians(335), 6)
    pts = [(cx + r * np.cos(a), cy - r * np.sin(a)) for a in angs]
    for i in range(len(pts) - 1):
        cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(pts[i][0]), Inches(pts[i][1]),
                                    Inches(pts[i + 1][0]), Inches(pts[i + 1][1]))
        cn.line.color.rgb = color; cn.line.width = Pt(w)
        if i == len(pts) - 2:
            ln = cn.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle"}))


def crop_thumb(src, dst, thr=20, pad=8):
    im = Image.open(src).convert("RGB")
    a = np.asarray(im).astype(int)
    bg = a[0, 0]
    mask = np.abs(a - bg).max(-1) > thr
    ys, xs = np.where(mask)
    if len(xs):
        box = (max(int(xs.min()) - pad, 0), max(int(ys.min()) - pad, 0),
               min(int(xs.max()) + pad, a.shape[1]), min(int(ys.max()) + pad, a.shape[0]))
        im = im.crop(box)
    im.save(dst)
    return dst


# ============================ SLIDE 1 — Title ============================
s = slide()
textbox(s, "Affordance-Grasp Conditioned\nSpill-Free Trajectory Optimization".split("\n"),
        0.9, 2.35, 11.5, 1.7, size=36, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
textbox(s, [[("Does ", {}), ("how you grasp", {"bold": True, "color": NAVY}),
             (" a cup change ", {}),
             ("how safely you can carry", {"bold": True, "color": NAVY}), (" it?", {})]],
        0.92, 4.15, 11.5, 0.6, size=20, color=SLATE)
textbox(s, "Jihoon Yun   ·   Motion Planning   ·   2026-06-08",
        0.92, 5.0, 11.5, 0.5, size=15, color=GRAY)

# ====================== SLIDE 2 — Big picture (native, editable) ========
s = slide()
title(s, "My thesis picks the grasp — this project plans the carry")

# region backgrounds
rrect(s, 0.30, 1.55, 7.05, 4.80, fill=THBG, line=None)
rrect(s, 7.55, 1.55, 5.45, 4.80, fill=PRBG, line=None)
textbox(s, [[("MY THESIS", {"bold": True, "color": NAVY, "size": 13}),
             ("  ·  Affordance & taxonomy-aware grasp pose generation",
              {"color": NAVY, "size": 12.5})]],
        0.45, 1.66, 6.8, 0.55, size=12.5, color=NAVY, align=PP_ALIGN.CENTER)
textbox(s, "THIS PROJECT  ·  spill-free carry (motion planning)",
        7.6, 1.66, 5.35, 0.55, size=12.5, color=DRED, bold=True, align=PP_ALIGN.CENTER)

# left: inputs -> grasp generation box
set_text(rrect(s, 0.55, 2.65, 1.7, 0.7, line=BLUE),
         [[("object", {})], [("mesh", {})]], size=11)
set_text(rrect(s, 0.55, 3.75, 2.1, 1.15, line=BLUE),
         [[("task (language)", {"size": 11})],
          [("\"pour\" / \"drink\" / \"hand over\"", {"size": 10, "color": SLATE})]], size=11)
set_text(rrect(s, 2.95, 3.0, 1.55, 1.3, line=BLUE),
         [[("Grasp", {"bold": True})], [("generation", {"bold": True})]], size=12.5)
arrow(s, 2.27, 2.95, 2.93, 3.4)
arrow(s, 2.67, 4.1, 2.93, 3.85)

# left: 3 grasp thumbnails (separate, editable images) + plain labels
thumbs = [("grasp_drinking.png", "drink", GREEN),
          ("grasp_pouring.png", "pour", BLUE),
          ("grasp_handover.png", "hand-over", PURPLE)]
ys = [2.45, 3.85, 5.25]
imx, imw = 4.95, 1.45
for (fn, lbl, col), yy in zip(thumbs, ys):
    cpath = crop_thumb(os.path.join(FIGS, fn), os.path.join(FIGS, fn[:-4] + "_c.png"))
    iw, ih = Image.open(cpath).size
    h = imw * ih / iw
    s.shapes.add_picture(cpath, Inches(imx), Inches(yy - h / 2), Inches(imw), Inches(h))
    arrow(s, 4.55, 3.65, imx - 0.02, yy, color=col, w=1.4)
    textbox(s, lbl, imx + imw + 0.05, yy - 0.2, 1.0, 0.4, size=11, color=col, bold=True,
            anchor=MSO_ANCHOR.MIDDLE)
textbox(s, "same mug, task → different grasp", 3.4, 6.05, 3.9, 0.35, size=10,
        color=SLATE, italic=True, align=PP_ALIGN.CENTER)

# cross arrow: a grasp -> this project
arrow(s, 6.55, 3.55, 7.75, 2.95, color=NAVY, w=2.5)
textbox(s, "a grasp", 6.55, 2.85, 1.2, 0.35, size=11, color=NAVY, bold=True)

# right: inputs -> optimization -> trajectory
set_text(rrect(s, 7.8, 2.45, 4.95, 0.7, line=DRED),
         "inputs:  grasp pose  +  start & goal EE poses", size=12)
arrow(s, 10.27, 3.15, 10.27, 3.48)
set_text(rrect(s, 8.25, 3.5, 4.05, 0.95, line=RED),
         [[("Spill-free trajectory optimization", {"bold": True})],
          [("(CHOMP + spill cost — ours)", {"size": 11})]], size=12.5)
arrow(s, 10.27, 4.45, 10.27, 4.78)
# trajectory output icon (editable shapes)
yb = 6.0
arc_curve(s, 8.85, 11.7, yb, 0.95, RED, w=2.6, n=16)
oval(s, 8.85, yb, 0.18, GREEN)
oval(s, 11.7, yb, 0.18, NAVY)
textbox(s, "start", 8.5, yb + 0.08, 0.9, 0.3, size=9.5, color=GREEN, align=PP_ALIGN.CENTER)
textbox(s, "goal", 11.35, yb + 0.08, 0.9, 0.3, size=9.5, color=NAVY, align=PP_ALIGN.CENTER)
textbox(s, "no-spill SE(3) carry trajectory", 7.8, 6.10, 4.95, 0.32, size=10.5,
        color=DRED, bold=True, align=PP_ALIGN.CENTER)

# bottom synthesis band: thesis + project merge into one outcome
arrow(s, 3.8, 6.40, 6.2, 6.52, color=NAVY, w=2.0)
arrow(s, 10.2, 6.40, 7.4, 6.52, color=NAVY, w=2.0)
band = rrect(s, 1.5, 6.55, 10.3, 0.62, fill=NAVY, line=None)
set_text(band, [[("⟹  Together: a motion-aware grasp-pose selector",
                  {"bold": True, "color": WHITE, "size": 14}),
                 ("   — pick the grasp the spill-free planner carries with the largest margin",
                  {"color": WHITE, "size": 11.5})]], align=PP_ALIGN.CENTER)
pagenum(s, 2)

# ============ SLIDE 3 — Problem & key insight (merged) ==================
s = slide()
t = title(s, "Problem — the danger is motion, not just tilt")
textbox(s, [[("Given a grasped mug + start/goal poses, plan a 6-DoF trajectory that is ", {}),
             ("spill-free, obstacle-avoiding, and smooth.", {"bold": True})]],
        0.7, t + 0.02, 12.0, 0.5, size=16, color=SLATE)
b = place_image(s, "spill_cone.png", t + 0.6, max_w=10.3, max_h=3.7)
textbox(s, [[("No-spill ⇔ apparent-up  a − g  stays inside the rim cone  θmax.   ",
              {"bold": True, "color": NAVY}),
             ("Static  tilt ≤ θmax  is the easy part; under acceleration  a − g  tilts away "
              "from vertical — that dynamic margin is the real constraint.", {})]],
        0.7, b + 0.12, 12.0, 0.8, size=14, color=SLATE)
pagenum(s, 3)

# ==================== SLIDE 4 — Why grasp matters =======================
s = slide()
t = title(s, "Why the grasp matters — it reshapes the spill margin")

# native mechanism diagram (each part editable): two grips, same arm motion
dv = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.62), Inches(1.7), Pt(1.2), Inches(2.3))
dv.fill.solid(); dv.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
dv.line.fill.background(); dv.shadow.inherit = False

# handle panel — large offset r
textbox(s, "handle grip  —  large offset r", 1.55, 1.6, 3.6, 0.4, size=12.5,
        color=GREEN, bold=True, align=PP_ALIGN.CENTER)
gripper(s, 2.15, 3.55)
textbox(s, "EE / wrist", 1.5, 3.72, 1.3, 0.3, size=9.5, color=SLATE, align=PP_ALIGN.CENTER)
cup_shape(s, 3.75, 2.0, 0.82, 1.0, GREEN)
arrow(s, 2.35, 3.5, 3.55, 2.95, color=GREEN, w=2.0, dashed=True)
textbox(s, "r", 2.88, 2.92, 0.3, 0.3, size=14, color=GREEN, bold=True)
arrow(s, 3.75, 2.45, 4.7, 2.45, color=RED, w=2.6)
textbox(s, "a_cup", 4.72, 2.28, 0.9, 0.3, size=11, color=RED)
rot_arc(s, 2.15, 3.28, 0.3, BLUE)
textbox(s, "α (same)", 0.95, 2.82, 1.1, 0.3, size=10, color=BLUE, align=PP_ALIGN.RIGHT)

# body panel — small offset r
textbox(s, "body grip  —  small offset r", 7.4, 1.6, 3.6, 0.4, size=12.5,
        color=RED, bold=True, align=PP_ALIGN.CENTER)
gripper(s, 8.7, 3.55)
textbox(s, "EE / wrist", 8.05, 3.72, 1.3, 0.3, size=9.5, color=SLATE, align=PP_ALIGN.CENTER)
cup_shape(s, 8.98, 2.35, 0.74, 0.9, RED)
arrow(s, 8.8, 3.5, 8.98, 3.25, color=RED, w=2.0, dashed=True)
textbox(s, "r", 8.48, 3.0, 0.3, 0.3, size=13, color=RED, bold=True)
arrow(s, 8.98, 2.82, 9.5, 2.82, color=RED, w=2.6)
textbox(s, "a_cup", 9.52, 2.65, 0.9, 0.3, size=11, color=RED)
rot_arc(s, 8.7, 3.28, 0.3, BLUE)
textbox(s, "α (same)", 7.5, 2.82, 1.1, 0.3, size=10, color=BLUE, align=PP_ALIGN.RIGHT)

# shared equation
textbox(s, [[("a = a_ee + α×r + ω×(ω×r)", {"bold": True, "color": NAVY, "size": 15}),
             ("      same arm motion, different offset r → different acceleration at the liquid",
              {"size": 12})]],
        0.7, 4.05, 12.0, 0.45, align=PP_ALIGN.CENTER)
b = 4.45
textbox(s, [
    [("•  Grasp reshapes the margin", {"bold": True, "color": NAVY}),
     (" — it sets the cup's offset r (lever effect on  a_cup) and how easily the arm "
      "keeps the cup upright (reachability).", {})],
    [("•  So which grasp is safest is ", {}),
     ("measured, not assumed", {"bold": True, "color": DRED}),
     (" — the planner scores each grasp  (→ next: Method).", {})],
], 0.7, b + 0.22, 12.1, 1.6, size=16.5, color=SLATE, line_spacing=1.25)
pagenum(s, 4)

# ====================== SLIDE 5 — Related work =========================
s = slide()
title(s, "Related work & the gap")
rows = [
    ["Line of work", "What they do", "Gap"],
    ["Anti-slosh / spill-free planning\n(Muchacho '22, Abderezaei '24)",
     "spill-free trajectory for a fixed container pose",
     "grasp is a fixed boundary condition"],
    ["Waiter / non-prehensile\n(Heins '23, Selvaggio '23)",
     "keep object from tipping / sliding (friction cone)",
     "rigid object, no liquid, no grasp choice"],
    ["Affordance / task-oriented grasping\n(Dang & Allen '12, Quispe & Stilman '16)",
     "pick a task-appropriate grasp",
     "stops at selection — never carries it"],
]
tbl_shape = s.shapes.add_table(4, 3, Inches(0.7), Inches(1.55),
                               Inches(11.9), Inches(3.2))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(4.1)
tbl.columns[1].width = Inches(4.2)
tbl.columns[2].width = Inches(3.6)
for ci in range(3):
    for ri in range(4):
        cell = tbl.cell(ri, ci)
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADBG if ri == 0 else (ROWBG if ri % 2 else WHITE)
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = rows[ri][ci]
        r.font.name = FONT
        r.font.size = Pt(14 if ri == 0 else 12.5)
        r.font.bold = (ri == 0)
        r.font.color.rgb = WHITE if ri == 0 else SLATE
textbox(s, [[("Nobody connects grasp choice to a physical spill margin. ",
              {"bold": True, "color": DRED, "size": 17}),
             ("We make the trajectory optimizer ", {"size": 16}),
             ("score grasps", {"bold": True, "size": 16}),
             (" by how safely they carry.", {"size": 16})]],
        0.7, 5.05, 12.0, 0.9, size=16, color=SLATE)
pagenum(s, 5)

# ====== SLIDE 6 — Method: spill-aware CHOMP, run per grasp (merged) ======
s = slide()
title(s, "Method — spill-aware CHOMP, run per grasp as an evaluator")

# cost formulation header (2 lines: objective + update, then spill-cost definition)
textbox(s, [
    [("J(ξ) = J_smooth + J_obs + λₛ J_spill",
      {"bold": True, "color": NAVY, "size": 16}),
     ("     (course CHOMP, Lec. 8  +  our spill term) ;     ξ ← ξ − η A⁻¹ ∇J",
      {"size": 12.5, "color": SLATE})],
    [("J_spill = Σₜ max( 0,  cos θmax − n·û )² ,    û = (a − g) / ‖a − g‖ ,    n = cup axis",
      {"bold": True, "color": DRED, "size": 13.5})],
], 0.5, 1.48, 12.3, 0.9, align=PP_ALIGN.CENTER, line_spacing=1.15)

# left: candidate grasps
textbox(s, "candidate grasps\n(taxonomy / region)", 0.45, 2.52, 3.0, 0.6, size=12,
        color=SLATE, bold=True, align=PP_ALIGN.CENTER)
cand = [("handle grip", GREEN, 3.05), ("rim grip", BLUE, 3.85), ("body grip", RED, 4.65)]
for name, col, yy in cand:
    set_text(rrect(s, 0.6, yy, 2.45, 0.58, line=col),
             [[(name, {"bold": True, "color": col})]], size=13)

# middle: per-grasp planner box
box = rrect(s, 3.95, 3.4, 2.35, 1.2, line=RED)
set_text(box, [[("Spill-aware CHOMP", {"bold": True})], [("(per grasp)", {"size": 11})]], size=13)
for _, col, yy in cand:
    arrow(s, 3.07, yy + 0.29, 3.93, 4.0, color=col, w=1.5)

# right: achievable-margin bars
textbox(s, "achievable margin to spill cone  (score)", 6.9, 2.57, 5.9, 0.4, size=12,
        color=NAVY, bold=True, align=PP_ALIGN.CENTER)
arrow(s, 6.32, 4.0, 7.45, 4.0, color=LGRAY, w=2.0)
bx0, maxw = 7.55, 4.55
for name, col, val, yy in [("handle", GREEN, 32, 3.05), ("rim", BLUE, 22, 3.85),
                           ("body", RED, 12, 4.65)]:
    w = maxw * val / 34.0
    bar(s, bx0, yy, w, 0.55, col)
    textbox(s, name, bx0 + 0.08, yy, 1.4, 0.55, size=12, color=WHITE, bold=True,
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, str(val), bx0 + w + 0.08, yy, 0.7, 0.55, size=13, color=col, bold=True,
            anchor=MSO_ANCHOR.MIDDLE)
textbox(s, "wider margin = easier, safer carry  (illustrative)", 7.5, 5.35, 5.2, 0.4,
        size=10, color=GRAY, italic=True)

# bottom: feedback loop (future / thesis)
arrow(s, 11.6, 5.85, 1.8, 5.85, color=NAVY, w=1.8, dashed=True)
textbox(s, [[("future (thesis): spill margin → grasp-generation reward — close the loop",
              {"italic": True, "color": NAVY, "size": 12})]],
        1.5, 6.03, 10.4, 0.4, align=PP_ALIGN.CENTER)
pagenum(s, 6)

# ===================== SLIDE 7 — Evaluation & plan =====================
s = slide()
t = title(s, "Evaluation & plan")
b = place_image(s, "ablation_plan.png", t + 0.05, max_w=11.5, max_h=3.0)
textbox(s, [
    [("• 3-way ablation: ", {"bold": True}),
     ("linear interp / CHOMP / CHOMP + spill — on a mug, with and without obstacles.", {})],
    [("• Metrics: ", {"bold": True}),
     ("spill ratio, path length, smoothness (jerk), plan time; run across grasps to show the "
      "margin differs.", {})],
    [("• Demo: ", {"bold": True}),
     ("SAPIEN particles on the liquid surface → count how many cross the rim.", {})],
    [("• Timeline: ", {"bold": True}),
     ("implementation + ablation within ~1 week  (proposal 06-08 → report 06-15).", {})],
], 0.7, b + 0.12, 12.1, 2.1, size=14.5, color=SLATE, line_spacing=1.12)
pagenum(s, 7)

out = os.path.join(DOCS, "slides_proposal_editable.pptx")
prs.save(out)
print("wrote", out, "(", len(prs.slides._sldIdLst), "slides )")
