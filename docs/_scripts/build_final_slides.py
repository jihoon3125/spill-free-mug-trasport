"""Build the 2-slide final deck (2-min talk): Problem + Results/Lessons.
Output: docs/slides_final.pptx   (editable in PowerPoint)
"""
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[2]
FIG = ROOT / "figs"
SW, SH = Inches(13.333), Inches(7.5)
DARK = RGBColor(0x20, 0x20, 0x20)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)

prs = Presentation()
prs.slide_width = SW; prs.slide_height = SH
blank = prs.slide_layouts[6]


def txt(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, space=6):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (s, sz, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        r = p.add_run(); r.text = s
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col
        r.font.name = "Arial"
    return tb


def pic(slide, path, l, t, w):
    im = Image.open(path); ar = im.height / im.width
    return slide.shapes.add_picture(str(path), Inches(l), Inches(t),
                                    width=Inches(w), height=Inches(w * ar))


def pic_h(slide, path, l, t, h):
    """Place by target height; returns width used (inches)."""
    im = Image.open(path); ar = im.width / im.height
    slide.shapes.add_picture(str(path), Inches(l), Inches(t),
                             height=Inches(h), width=Inches(h * ar))
    return h * ar


# ---------------- Slide 1 — Problem ----------------
s = prs.slides.add_slide(blank)
txt(s, 0.5, 0.22, 12.3, 0.9,
    [("Grasp-Conditioned Spill-Free Mug Transport", 30, True, DARK)])
txt(s, 0.5, 1.0, 12.3, 0.4,
    [("Motion Planning Term Project  ·  Jihoon Yoon", 14, False, GRAY)])
txt(s, 0.5, 1.45, 12.3, 1.75, [
    ("Given an affordance grasp (from my thesis), generate a 6-DoF arm trajectory "
     "that carries a mug WITHOUT spilling.", 15, False, DARK),
    ("Key physics — the “waiter problem”: even an upright cup spills under acceleration, "
     "because the water feels effective gravity g − a.", 15, False, DARK),
    ("Approach: CHOMP + a differentiable spill cost (apparent-gravity must stay in the "
     "cup’s rim cone); evaluate per grasp.", 15, False, DARK),
])
gw = pic_h(s, FIG / "_gallery_still.png", 0, 3.3, 3.75)
# center horizontally
prs_w_in = 13.333
s.shapes[-1].left = Inches((prs_w_in - gw) / 2)
txt(s, 0.5, 7.12, 12.3, 0.35,
    [("Same tasks, 3 scenarios — top: Min-jerk spills · bottom: Ours keeps it in "
      "(illustration of the reduced slosh model)", 11, False, GRAY)],
    align=PP_ALIGN.CENTER, space=0)

# ---------------- Slide 2 — Results & Lessons ----------------
s = prs.slides.add_slide(blank)
txt(s, 0.5, 0.18, 12.3, 0.6, [("Results & Key Lessons", 28, True, DARK)])

# 2x2 grid of result figures
H = 1.75
quad = [(0.7, 0.78, "fig_method_comparison.png", "Ours: lowest spill + feasible"),
        (7.0, 0.78, "fig_gamma_pareto.png", "gamma trades spill vs feasibility"),
        (0.7, 2.95, "fig_filllevel.png", "robust up to ~half-full"),
        (7.0, 2.95, "fig_tstar_offset.png", "T* varies by grasp, not from geometry")]
for x, top, f, cap in quad:
    w = pic_h(s, FIG / f, x, top, H)
    txt(s, x, top + H + 0.01, w, 0.3, [(cap, 10, True, GREEN)],
        align=PP_ALIGN.CENTER, space=0)

txt(s, 0.4, 4.95, 12.6, 2.4, [
    ("Ours is the only method that is both nearly spill-free (1.1%) and joint-limit "
     "feasible (94%); baselines either spill or violate joint limits.", 13, False, DARK),
    ("The spill weight gamma trades spill vs feasibility — gamma≈1 is the sweet spot "
     "(larger gamma triggers a degenerate, infeasible high-acceleration mode).", 13, False, DARK),
    ("Robust up to a ~half-full mug (95% of grasps kept spill-free), degrading gracefully "
     "as it nears the brim.", 13, False, DARK),
    ("Fastest feasible spill-free time varies by grasp (0.7–1.6 s), not predicted by "
     "geometry (r≈0); 36% of grasps unreachable ⇒ measure, don’t assume.", 13, True, GREEN),
])

out = ROOT / "docs/slides_final.pptx"
prs.save(str(out))
print(f"saved {out}")
